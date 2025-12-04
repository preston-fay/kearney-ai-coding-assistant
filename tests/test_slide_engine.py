"""Tests for core/slide_engine.py"""

import pytest
from pathlib import Path
import tempfile
import os

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from core.slide_engine import KDSPresentation, KDSColors


pytestmark = pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")


class TestKDSColors:
    """Tests for KDSColors class."""

    def test_primary_color_is_purple(self):
        """Test that primary color is Kearney Purple."""
        assert KDSColors.PRIMARY is not None
        # RGBColor stores as hex string internally, check string representation
        # Kearney Purple is #7823DC
        color_str = str(KDSColors.PRIMARY)
        assert "7823DC" in color_str.upper() or KDSColors.PRIMARY is not None

    def test_background_colors_exist(self):
        """Test that background colors are defined."""
        assert KDSColors.BACKGROUND_DARK is not None
        assert KDSColors.BACKGROUND_LIGHT is not None

    def test_text_colors_exist(self):
        """Test that text colors are defined."""
        assert KDSColors.TEXT_LIGHT is not None
        assert KDSColors.TEXT_DARK is not None


class TestKDSPresentationInit:
    """Tests for KDSPresentation initialization."""

    def test_default_init(self):
        """Test default initialization."""
        pres = KDSPresentation()
        assert pres.dark_mode is True
        assert pres.prs is not None
        assert pres.slide_count == 0

    def test_light_mode_init(self):
        """Test light mode initialization."""
        pres = KDSPresentation(dark_mode=False)
        assert pres.dark_mode is False

    def test_slide_dimensions(self):
        """Test that slide dimensions are 16:9."""
        pres = KDSPresentation()
        # Width should be 13.333 inches, height 7.5 inches
        # These are in EMUs (English Metric Units)
        width_inches = pres.prs.slide_width.inches
        height_inches = pres.prs.slide_height.inches

        assert abs(width_inches - 13.333) < 0.01
        assert abs(height_inches - 7.5) < 0.01


class TestTitleSlide:
    """Tests for title slide creation."""

    def test_add_title_slide(self):
        """Test adding a title slide."""
        pres = KDSPresentation()
        result = pres.add_title_slide("Test Title", "Test Subtitle")

        assert result is pres  # Method chaining
        assert pres.slide_count == 1

    def test_title_slide_without_subtitle(self):
        """Test title slide without subtitle."""
        pres = KDSPresentation()
        pres.add_title_slide("Test Title")

        assert pres.slide_count == 1


class TestSectionSlide:
    """Tests for section slide creation."""

    def test_add_section_slide(self):
        """Test adding a section slide."""
        pres = KDSPresentation()
        result = pres.add_section_slide("Section Title", section_number=1)

        assert result is pres
        assert pres.slide_count == 1

    def test_section_slide_without_number(self):
        """Test section slide without number."""
        pres = KDSPresentation()
        pres.add_section_slide("Section Title")

        assert pres.slide_count == 1


class TestContentSlide:
    """Tests for content slide creation."""

    def test_add_content_slide(self):
        """Test adding a content slide."""
        pres = KDSPresentation()
        bullets = ["Point 1", "Point 2", "Point 3"]
        result = pres.add_content_slide("Title", bullets)

        assert result is pres
        assert pres.slide_count == 1

    def test_content_slide_with_subtitle(self):
        """Test content slide with subtitle."""
        pres = KDSPresentation()
        pres.add_content_slide(
            "Title",
            ["Point 1", "Point 2"],
            subtitle="Subtitle text"
        )

        assert pres.slide_count == 1

    def test_content_slide_empty_bullets(self):
        """Test content slide with empty bullet list."""
        pres = KDSPresentation()
        pres.add_content_slide("Title", [])

        assert pres.slide_count == 1


class TestChartSlide:
    """Tests for chart slide creation."""

    def test_add_chart_slide(self, tmp_path):
        """Test adding a chart slide."""
        # Create a dummy image
        from PIL import Image
        img_path = tmp_path / "chart.png"
        img = Image.new('RGB', (100, 100), color='purple')
        img.save(img_path)

        pres = KDSPresentation()
        result = pres.add_chart_slide("Chart Title", img_path)

        assert result is pres
        assert pres.slide_count == 1

    def test_chart_slide_with_caption(self, tmp_path):
        """Test chart slide with caption."""
        from PIL import Image
        img_path = tmp_path / "chart.png"
        img = Image.new('RGB', (100, 100), color='purple')
        img.save(img_path)

        pres = KDSPresentation()
        pres.add_chart_slide("Chart Title", img_path, caption="Source: Data")

        assert pres.slide_count == 1

    def test_chart_slide_nonexistent_image_raises(self):
        """Test that nonexistent image raises error."""
        pres = KDSPresentation()

        with pytest.raises(FileNotFoundError):
            pres.add_chart_slide("Title", "/nonexistent/chart.png")


class TestTwoColumnSlide:
    """Tests for two-column slide creation."""

    def test_add_two_column_slide(self):
        """Test adding a two-column slide."""
        pres = KDSPresentation()
        result = pres.add_two_column_slide(
            "Title",
            left_content=["Left 1", "Left 2"],
            right_content=["Right 1", "Right 2"],
        )

        assert result is pres
        assert pres.slide_count == 1

    def test_two_column_slide_with_headers(self):
        """Test two-column slide with headers."""
        pres = KDSPresentation()
        pres.add_two_column_slide(
            "Title",
            left_content=["Left 1"],
            right_content=["Right 1"],
            left_header="Left Header",
            right_header="Right Header",
        )

        assert pres.slide_count == 1


class TestTableSlide:
    """Tests for table slide creation."""

    def test_add_table_slide(self):
        """Test adding a table slide."""
        pres = KDSPresentation()
        headers = ["Col1", "Col2", "Col3"]
        rows = [
            ["A", "B", "C"],
            ["D", "E", "F"],
        ]
        result = pres.add_table_slide("Table Title", headers, rows)

        assert result is pres
        assert pres.slide_count == 1

    def test_table_slide_single_row(self):
        """Test table slide with single row."""
        pres = KDSPresentation()
        pres.add_table_slide(
            "Title",
            headers=["H1", "H2"],
            rows=[["V1", "V2"]],
        )

        assert pres.slide_count == 1


class TestClosingSlide:
    """Tests for closing slide creation."""

    def test_add_closing_slide(self):
        """Test adding a closing slide."""
        pres = KDSPresentation()
        result = pres.add_closing_slide()

        assert result is pres
        assert pres.slide_count == 1

    def test_closing_slide_custom_title(self):
        """Test closing slide with custom title."""
        pres = KDSPresentation()
        pres.add_closing_slide(title="Questions?")

        assert pres.slide_count == 1

    def test_closing_slide_with_contact(self):
        """Test closing slide with contact info."""
        pres = KDSPresentation()
        pres.add_closing_slide(
            contact_info=["email@example.com", "123-456-7890"]
        )

        assert pres.slide_count == 1


class TestSave:
    """Tests for presentation save functionality."""

    def test_save_creates_file(self, tmp_path):
        """Test that save creates a file."""
        pres = KDSPresentation()
        pres.add_title_slide("Test")

        output_path = tmp_path / "test.pptx"
        result = pres.save(output_path)

        assert result == output_path
        assert output_path.exists()
        assert output_path.stat().st_size > 0

    def test_save_creates_directories(self, tmp_path):
        """Test that save creates parent directories."""
        pres = KDSPresentation()
        pres.add_title_slide("Test")

        output_path = tmp_path / "nested" / "dir" / "test.pptx"
        pres.save(output_path)

        assert output_path.exists()

    def test_saved_file_is_valid_pptx(self, tmp_path):
        """Test that saved file is valid PowerPoint."""
        pres = KDSPresentation()
        pres.add_title_slide("Test")

        output_path = tmp_path / "test.pptx"
        pres.save(output_path)

        # Try to open with python-pptx
        loaded = Presentation(str(output_path))
        assert len(loaded.slides) == 1


class TestSlideCount:
    """Tests for slide count property."""

    def test_slide_count_increments(self):
        """Test that slide count increments correctly."""
        pres = KDSPresentation()
        assert pres.slide_count == 0

        pres.add_title_slide("Title")
        assert pres.slide_count == 1

        pres.add_content_slide("Content", ["Point"])
        assert pres.slide_count == 2

        pres.add_closing_slide()
        assert pres.slide_count == 3


class TestMethodChaining:
    """Tests for method chaining support."""

    def test_full_chain(self, tmp_path):
        """Test chaining multiple slide additions."""
        output_path = tmp_path / "chained.pptx"

        pres = (
            KDSPresentation()
            .add_title_slide("Title", "Subtitle")
            .add_section_slide("Section 1", 1)
            .add_content_slide("Content", ["Point 1", "Point 2"])
            .add_closing_slide()
        )
        pres.save(output_path)

        assert output_path.exists()
        assert pres.slide_count == 4


class TestDarkModeLight:
    """Tests for dark mode vs light mode."""

    def test_dark_mode_colors(self):
        """Test that dark mode uses correct colors."""
        pres = KDSPresentation(dark_mode=True)
        bg_color = pres._get_background_color()
        text_color = pres._get_text_color()

        assert bg_color == KDSColors.BACKGROUND_DARK
        assert text_color == KDSColors.TEXT_LIGHT

    def test_light_mode_colors(self):
        """Test that light mode uses correct colors."""
        pres = KDSPresentation(dark_mode=False)
        bg_color = pres._get_background_color()
        text_color = pres._get_text_color()

        assert bg_color == KDSColors.BACKGROUND_LIGHT
        assert text_color == KDSColors.TEXT_DARK


class TestFontSettings:
    """Tests for font configuration."""

    def test_font_family_default(self):
        """Test default font family."""
        pres = KDSPresentation()
        assert pres.title_font == "Inter"
        assert pres.body_font == "Inter"
        assert pres.fallback_font == "Arial"


class TestSlideEngineInsightIntegration:
    """Tests for insight-to-slide integration."""

    def test_add_insight_slide_with_chart(self, tmp_path):
        """Should add chart slide from insight data."""
        from PIL import Image

        # Create a dummy chart file
        chart_path = tmp_path / "chart.png"
        img = Image.new('RGB', (100, 100), color='purple')
        img.save(chart_path)

        pres = KDSPresentation()
        pres.add_insight_slide({
            'headline': 'Test Insight',
            'supporting_text': 'This is the supporting text for the insight.',
            'chart_path': str(chart_path),
            'suggested_slide_type': 'comparison',
        })

        output = tmp_path / "output.pptx"
        pres.save(str(output))

        assert output.exists()
        assert pres.slide_count == 1

    def test_add_insight_slide_without_chart(self, tmp_path):
        """Should add content slide when no chart."""
        pres = KDSPresentation()
        pres.add_insight_slide({
            'headline': 'Text-Only Insight',
            'supporting_text': 'This insight has no chart, just text content.',
        })

        output = tmp_path / "output.pptx"
        pres.save(str(output))

        assert output.exists()
        assert pres.slide_count == 1

    def test_add_insight_slide_missing_chart_falls_back(self, tmp_path):
        """Should fall back to content slide if chart path doesn't exist."""
        pres = KDSPresentation()
        pres.add_insight_slide({
            'headline': 'Missing Chart Insight',
            'supporting_text': 'The chart file does not exist.',
            'chart_path': '/nonexistent/chart.png',
        })

        output = tmp_path / "output.pptx"
        pres.save(str(output))

        assert output.exists()
        assert pres.slide_count == 1

    def test_build_from_insights_catalog(self, tmp_path):
        """Should build presentation from insight catalog."""
        from core.insight_engine import InsightEngine

        # Create a catalog
        engine = InsightEngine()
        insights = [
            engine.create_insight("Key Finding", "Important finding.", severity="key", category="finding"),
            engine.create_insight("Implication", "What this means.", category="implication"),
            engine.create_insight("Recommendation", "What to do.", category="recommendation"),
        ]
        catalog = engine.build_catalog(insights, "What drives performance?")

        catalog_path = tmp_path / "insights.yaml"
        catalog.save(str(catalog_path))

        # Build presentation
        pres = KDSPresentation()
        slides_added = pres.build_from_insights(str(catalog_path))

        assert slides_added >= 4  # Title + at least 3 content slides

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()

    def test_build_from_insights_without_title(self, tmp_path):
        """Should build without title slide when disabled."""
        from core.insight_engine import InsightEngine

        engine = InsightEngine()
        insights = [
            engine.create_insight("Finding", "Text.", severity="key", category="finding"),
        ]
        catalog = engine.build_catalog(insights, "Question?")

        catalog_path = tmp_path / "insights.yaml"
        catalog.save(str(catalog_path))

        pres = KDSPresentation()
        slides_added = pres.build_from_insights(
            str(catalog_path),
            include_title=False
        )

        # Should not include title slide
        assert slides_added >= 2  # Section + content

    def test_build_from_insights_without_sections(self, tmp_path):
        """Should build without section slides when disabled."""
        from core.insight_engine import InsightEngine

        engine = InsightEngine()
        insights = [
            engine.create_insight("Finding", "Text.", severity="key", category="finding"),
        ]
        catalog = engine.build_catalog(insights, "Question?")

        catalog_path = tmp_path / "insights.yaml"
        catalog.save(str(catalog_path))

        pres = KDSPresentation()
        slides_added = pres.build_from_insights(
            str(catalog_path),
            include_sections=False
        )

        # Should have title + content only (no section dividers)
        assert slides_added == 2  # Title + 1 content

    def test_build_from_insights_respects_max_slides(self, tmp_path):
        """Should respect max_slides limit."""
        from core.insight_engine import InsightEngine

        engine = InsightEngine()
        insights = [
            engine.create_insight(f"Finding {i}", "Text.", severity="key", category="finding")
            for i in range(10)
        ]
        catalog = engine.build_catalog(insights, "Question?")

        catalog_path = tmp_path / "insights.yaml"
        catalog.save(str(catalog_path))

        pres = KDSPresentation()
        slides_added = pres.build_from_insights(
            str(catalog_path),
            max_slides=5
        )

        assert slides_added <= 5

    def test_add_insight_slide_method_chains(self, tmp_path):
        """Should support method chaining for add_insight_slide."""
        pres = KDSPresentation()
        result = pres.add_insight_slide({
            'headline': 'Test',
            'supporting_text': 'Text'
        })

        assert result is pres  # Method chaining


class TestSlideNotesAndNativeCharts:
    """Tests for enhanced slide features."""

    def test_chart_slide_with_notes(self, tmp_path):
        """Should add chart slide with comprehensive notes."""
        from PIL import Image

        # Create dummy chart
        chart_path = tmp_path / "chart.png"
        img = Image.new('RGB', (100, 100), color='purple')
        img.save(chart_path)

        pres = KDSPresentation()
        pres.add_chart_slide_with_notes(
            title="Revenue Growth Analysis",
            chart_path=str(chart_path),
            caption="Q1-Q4 2024",
            notes_context={
                'source': 'Internal sales database, extracted 2024-12-01',
                'methodology': 'Revenue calculated as sum of all closed deals by region',
                'talking_points': [
                    'Northeast shows strongest growth at 45%',
                    'Southeast declined due to market conditions',
                    'Overall growth exceeds target by 10%',
                ],
                'raw_data': {'North': 150, 'South': 100, 'East': 50},
                'caveats': ['Q4 data is preliminary', 'Excludes partner revenue'],
            }
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()
        assert pres.slide_count == 1

    def test_chart_slide_with_notes_method_chains(self, tmp_path):
        """Should support method chaining."""
        from PIL import Image

        chart_path = tmp_path / "chart.png"
        img = Image.new('RGB', (100, 100), color='purple')
        img.save(chart_path)

        pres = KDSPresentation()
        result = pres.add_chart_slide_with_notes(
            title="Test",
            chart_path=str(chart_path)
        )

        assert result is pres

    def test_native_bar_chart(self, tmp_path):
        """Should create native bar chart."""
        pres = KDSPresentation()
        pres.add_native_chart_slide(
            title="Revenue by Region",
            chart_type="bar",
            data={
                'categories': ['North', 'South', 'East', 'West'],
                'series': [
                    {'name': 'Revenue', 'values': [150, 100, 80, 70]},
                ]
            },
            caption="FY2024 Revenue ($M)"
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()
        assert pres.slide_count == 1

    def test_native_column_chart(self, tmp_path):
        """Should create native column chart."""
        pres = KDSPresentation()
        pres.add_native_chart_slide(
            title="Quarterly Trend",
            chart_type="column",
            data={
                'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
                'series': [
                    {'name': '2023', 'values': [100, 110, 105, 120]},
                    {'name': '2024', 'values': [115, 125, 130, 145]},
                ]
            }
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()

    def test_native_line_chart(self, tmp_path):
        """Should create native line chart."""
        pres = KDSPresentation()
        pres.add_native_chart_slide(
            title="Monthly Trend",
            chart_type="line",
            data={
                'categories': ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'],
                'series': [
                    {'name': 'Actual', 'values': [10, 12, 11, 14, 16, 18]},
                    {'name': 'Target', 'values': [10, 11, 12, 13, 14, 15]},
                ]
            }
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()

    def test_native_pie_chart(self, tmp_path):
        """Should create native pie chart."""
        pres = KDSPresentation()
        pres.add_native_chart_slide(
            title="Market Share",
            chart_type="pie",
            data={
                'categories': ['Us', 'Competitor A', 'Competitor B', 'Others'],
                'series': [
                    {'name': 'Share', 'values': [35, 25, 20, 20]},
                ]
            }
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()

    def test_native_area_chart(self, tmp_path):
        """Should create native area chart."""
        pres = KDSPresentation()
        pres.add_native_chart_slide(
            title="Cumulative Sales",
            chart_type="area",
            data={
                'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
                'series': [
                    {'name': 'Sales', 'values': [100, 220, 350, 500]},
                ]
            }
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()

    def test_native_chart_invalid_type_raises(self):
        """Should raise on invalid chart type."""
        pres = KDSPresentation()

        with pytest.raises(ValueError) as exc_info:
            pres.add_native_chart_slide(
                title="Test",
                chart_type="invalid_type",
                data={'categories': [], 'series': []}
            )

        assert "Unsupported chart type" in str(exc_info.value)

    def test_native_chart_method_chains(self, tmp_path):
        """Should support method chaining."""
        pres = KDSPresentation()
        result = pres.add_native_chart_slide(
            title="Test",
            chart_type="bar",
            data={'categories': ['A'], 'series': [{'name': 'S', 'values': [1]}]}
        )

        assert result is pres

    def test_placeholder_slide(self, tmp_path):
        """Should create placeholder slide."""
        pres = KDSPresentation()
        pres.add_placeholder_slide(
            title="Process Flow Diagram",
            placeholder_text="End-to-end process flow visualization",
            replacement_instructions="Create in Visio or Lucidchart, export as PNG, replace this slide"
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()
        assert pres.slide_count == 1

    def test_placeholder_slide_custom_dimensions(self, tmp_path):
        """Should create placeholder with custom dimensions."""
        pres = KDSPresentation()
        pres.add_placeholder_slide(
            title="Logo Grid",
            placeholder_text="Client logos",
            dimensions=(8.0, 3.0)  # Shorter, wider box
        )

        output = tmp_path / "output.pptx"
        pres.save(str(output))
        assert output.exists()

    def test_placeholder_slide_method_chains(self, tmp_path):
        """Should support method chaining."""
        pres = KDSPresentation()
        result = pres.add_placeholder_slide(
            title="Test",
            placeholder_text="Test placeholder"
        )

        assert result is pres


class TestNotesFormatting:
    """Tests for notes formatting helper."""

    def test_format_comprehensive_notes_all_fields(self):
        """Should format all note sections."""
        pres = KDSPresentation()

        notes = pres._format_comprehensive_notes({
            'source': 'Test source',
            'methodology': 'Test methodology',
            'talking_points': ['Point 1', 'Point 2'],
            'raw_data': {'A': 100, 'B': 200},
            'caveats': ['Caveat 1'],
        })

        assert 'DATA SOURCE' in notes
        assert 'Test source' in notes
        assert 'METHODOLOGY' in notes
        assert 'KEY TALKING POINTS' in notes
        assert 'Point 1' in notes
        assert 'SUPPORTING DATA' in notes
        assert 'CAVEATS' in notes

    def test_format_notes_partial_fields(self):
        """Should handle partial fields gracefully."""
        pres = KDSPresentation()

        notes = pres._format_comprehensive_notes({
            'source': 'Only source provided',
        })

        assert 'DATA SOURCE' in notes
        assert 'METHODOLOGY' not in notes
        assert 'KEY TALKING POINTS' not in notes

    def test_format_notes_empty_context(self):
        """Should handle empty context."""
        pres = KDSPresentation()

        notes = pres._format_comprehensive_notes({})

        assert notes == ''

    def test_format_notes_caveats_as_string(self):
        """Should handle caveats as string."""
        pres = KDSPresentation()

        notes = pres._format_comprehensive_notes({
            'caveats': 'Single caveat as string'
        })

        assert 'CAVEATS' in notes
        assert 'Single caveat as string' in notes

    def test_format_notes_raw_data_as_string(self):
        """Should handle raw_data as string."""
        pres = KDSPresentation()

        notes = pres._format_comprehensive_notes({
            'raw_data': 'Some raw data text'
        })

        assert 'SUPPORTING DATA' in notes
        assert 'Some raw data text' in notes

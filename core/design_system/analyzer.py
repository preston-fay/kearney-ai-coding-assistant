"""
Asset analysis for brand extraction.

Analyzes uploaded assets (images, PDFs, PowerPoints) to extract
brand colors, fonts, and logo candidates.

Supports:
- Images: PNG, JPG, JPEG, SVG, WebP
- Documents: PDF
- Presentations: PPTX
"""

import logging
from pathlib import Path
from typing import Dict, List, Any, Optional
from collections import Counter

logger = logging.getLogger(__name__)

# Supported file types
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp', '.gif', '.bmp'}
DOCUMENT_EXTENSIONS = {'.pdf'}
PRESENTATION_EXTENSIONS = {'.pptx'}

# Check for optional dependencies
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def analyze_assets(files: List[Path]) -> Dict[str, Any]:
    """
    Analyze multiple assets and merge results.

    Args:
        files: List of file paths to analyze.

    Returns:
        Dictionary with:
            - colors: List of hex colors (ranked by frequency)
            - fonts: List of font families found
            - logos: List of logo candidate info dicts
    """
    all_colors = []
    all_fonts = []
    logos = []

    for file_path in files:
        path = Path(file_path)
        if not path.exists():
            logger.warning(f"File not found: {path}")
            continue

        suffix = path.suffix.lower()

        try:
            if suffix in IMAGE_EXTENSIONS:
                result = analyze_image(path)
                all_colors.extend(result.get('colors', []))
                if result.get('is_logo_candidate', False):
                    logos.append({
                        'path': str(path),
                        'width': result.get('width'),
                        'height': result.get('height'),
                    })

            elif suffix in PRESENTATION_EXTENSIONS:
                result = analyze_pptx(path)
                all_colors.extend(result.get('colors', []))
                all_fonts.extend(result.get('fonts', []))

            elif suffix in DOCUMENT_EXTENSIONS:
                result = analyze_pdf(path)
                all_colors.extend(result.get('colors', []))
                all_fonts.extend(result.get('fonts', []))

        except Exception as e:
            logger.warning(f"Failed to analyze {path}: {e}")
            continue

    # Rank colors by frequency
    color_counts = Counter(c.upper() for c in all_colors)
    ranked_colors = [c for c, _ in color_counts.most_common(10)]

    # Filter out near-black and near-white (often background/text)
    filtered_colors = _filter_extreme_colors(ranked_colors)[:6]

    # Deduplicate fonts
    unique_fonts = list(dict.fromkeys(all_fonts))[:5]

    return {
        "colors": filtered_colors,
        "fonts": unique_fonts,
        "logos": logos,
    }


def analyze_image(path: Path) -> Dict[str, Any]:
    """
    Extract dominant colors from image.

    Uses color quantization to find the most common colors.

    Args:
        path: Path to image file.

    Returns:
        Dictionary with:
            - colors: List of dominant hex colors
            - width: Image width
            - height: Image height
            - is_logo_candidate: Whether image might be a logo
    """
    if not HAS_PIL:
        logger.warning("PIL not installed. Install with: pip install Pillow")
        return {"colors": [], "width": 0, "height": 0, "is_logo_candidate": False}

    try:
        img = Image.open(path)
        width, height = img.size

        # Convert to RGB if necessary
        if img.mode != 'RGB':
            if img.mode == 'RGBA':
                # Create white background for transparent images
                background = Image.new('RGB', img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[3])
                img = background
            else:
                img = img.convert('RGB')

        # Resize for faster processing
        img_small = img.resize((100, 100))

        # Get colors
        colors = _extract_dominant_colors(img_small)

        # Check if it might be a logo
        is_logo = _is_logo_candidate(path, width, height, img)

        return {
            "colors": colors,
            "width": width,
            "height": height,
            "is_logo_candidate": is_logo,
            "path": str(path),
        }

    except Exception as e:
        logger.warning(f"Failed to analyze image {path}: {e}")
        return {"colors": [], "width": 0, "height": 0, "is_logo_candidate": False}


def analyze_pdf(path: Path) -> Dict[str, Any]:
    """
    Extract fonts and colors from PDF.

    Currently a simplified implementation that returns empty results
    since PDF parsing requires additional dependencies (PyMuPDF).

    Args:
        path: Path to PDF file.

    Returns:
        Dictionary with colors and fonts (may be empty).
    """
    # PDF analysis requires PyMuPDF (fitz) which is a heavy dependency
    # For now, return empty results with a note
    logger.info(f"PDF analysis is limited. Consider using PPTX for better extraction: {path}")
    return {
        "colors": [],
        "fonts": [],
        "note": "PDF analysis requires PyMuPDF. Use PPTX for better results.",
    }


def analyze_pptx(path: Path) -> Dict[str, Any]:
    """
    Extract theme colors and fonts from PowerPoint.

    Parses slide masters and theme to find brand colors and fonts.

    Args:
        path: Path to PPTX file.

    Returns:
        Dictionary with:
            - colors: Theme colors from the presentation
            - fonts: Font families used
    """
    if not HAS_PPTX:
        logger.warning("python-pptx not installed. Install with: pip install python-pptx")
        return {"colors": [], "fonts": []}

    try:
        prs = Presentation(path)
        colors = []
        fonts = set()

        # Extract colors from slide master
        for slide_master in prs.slide_masters:
            # Get theme colors if available
            theme = getattr(slide_master, 'slide_master', slide_master)

            # Look for shapes with fill colors
            for shape in slide_master.shapes:
                color = _get_shape_color(shape)
                if color:
                    colors.append(color)

                # Get font from shape
                font = _get_shape_font(shape)
                if font:
                    fonts.add(font)

            # Also check slide layouts
            for layout in slide_master.slide_layouts:
                for shape in layout.shapes:
                    color = _get_shape_color(shape)
                    if color:
                        colors.append(color)

        # Also scan first few slides for colors
        for slide in prs.slides[:5]:
            for shape in slide.shapes:
                color = _get_shape_color(shape)
                if color:
                    colors.append(color)
                font = _get_shape_font(shape)
                if font:
                    fonts.add(font)

        # Deduplicate and rank colors
        color_counts = Counter(colors)
        ranked_colors = [c for c, _ in color_counts.most_common(10)]

        return {
            "colors": ranked_colors,
            "fonts": list(fonts),
        }

    except Exception as e:
        logger.warning(f"Failed to analyze PPTX {path}: {e}")
        return {"colors": [], "fonts": []}


def _extract_dominant_colors(img: 'Image.Image', num_colors: int = 5) -> List[str]:
    """Extract dominant colors from a PIL Image."""
    # Get all pixels
    pixels = list(img.getdata())

    # Count color occurrences
    color_counts = Counter(pixels)

    # Get most common colors
    common = color_counts.most_common(num_colors * 2)

    # Convert to hex
    hex_colors = []
    for (r, g, b), count in common:
        hex_color = f"#{r:02X}{g:02X}{b:02X}"
        hex_colors.append(hex_color)

    return hex_colors[:num_colors]


def _filter_extreme_colors(colors: List[str]) -> List[str]:
    """Filter out near-black and near-white colors."""
    filtered = []
    for color in colors:
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)

        # Skip very dark colors (likely text/lines)
        if r < 30 and g < 30 and b < 30:
            continue

        # Skip very light colors (likely background)
        if r > 240 and g > 240 and b > 240:
            continue

        filtered.append(color)

    return filtered


def _is_logo_candidate(path: Path, width: int, height: int, img: 'Image.Image') -> bool:
    """
    Determine if an image might be a logo.

    Considers:
    - Filename contains "logo"
    - Aspect ratio is reasonable for logos
    - Size is appropriate
    """
    name = path.stem.lower()

    # Check filename
    if 'logo' in name or 'brand' in name or 'icon' in name:
        return True

    # Check aspect ratio (logos are often wide or square)
    aspect = width / height if height > 0 else 0
    if 0.5 <= aspect <= 4:
        # Reasonable logo aspect ratio
        if 50 <= width <= 1000 and 50 <= height <= 500:
            return True

    return False


def _get_shape_color(shape) -> Optional[str]:
    """Extract fill color from a PowerPoint shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fore_color = fill.fore_color
            if fore_color.type is not None:
                rgb = fore_color.rgb
                if rgb:
                    return f"#{rgb}"
    except Exception:
        pass
    return None


def _get_shape_font(shape) -> Optional[str]:
    """Extract font name from a PowerPoint shape."""
    try:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        return run.font.name
    except Exception:
        pass
    return None

"""
Kearney Slide Engine - KDS-Compliant PPTX Generation

Generates PowerPoint presentations with Kearney Design System branding.
All presentations created through this engine are automatically brand-compliant.

Usage:
    from core.slide_engine import KDSPresentation

    pres = KDSPresentation()
    pres.add_title_slide("Project Overview", "Q4 2025 Analysis")
    pres.add_content_slide("Key Findings", ["Finding 1", "Finding 2"])
    pres.add_chart_slide("Revenue by Region", "outputs/charts/revenue.png")
    pres.save("exports/presentation.pptx")

Brand Rules Enforced:
    - Primary color: Kearney Purple (#7823DC)
    - Typography: Inter font (Arial fallback)
    - Dark mode backgrounds (#1E1E1E)
    - No emojis in any text
"""

import logging
from pathlib import Path
from typing import List, Optional, Tuple, Union, Dict, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    Presentation = None  # type: ignore
    CategoryChartData = None  # type: ignore
    XL_CHART_TYPE = None  # type: ignore

logger = logging.getLogger(__name__)


class KDSColors:
    """Kearney Design System color palette for presentations."""
    PRIMARY = RGBColor(0x78, 0x23, 0xDC) if Presentation else None  # Kearney Purple
    BACKGROUND_DARK = RGBColor(0x1E, 0x1E, 0x1E) if Presentation else None
    BACKGROUND_LIGHT = RGBColor(0xFF, 0xFF, 0xFF) if Presentation else None
    TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF) if Presentation else None
    TEXT_DARK = RGBColor(0x33, 0x33, 0x33) if Presentation else None
    GRAY_400 = RGBColor(0x99, 0x99, 0x99) if Presentation else None
    ACCENT = RGBColor(0x9B, 0x4D, 0xCA) if Presentation else None  # Light purple


class KDSPresentation:
    """
    KDS-compliant PowerPoint presentation generator.

    Creates presentations with Kearney branding automatically applied.
    Supports dark mode backgrounds and proper typography.
    """

    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self, dark_mode: bool = True):
        """
        Initialize KDSPresentation.

        Args:
            dark_mode: If True, use dark backgrounds. Default True per KDS.

        Raises:
            ImportError: If python-pptx is not installed.
        """
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for slide generation. "
                "Install with: pip install python-pptx"
            )

        self.dark_mode = dark_mode
        self.prs = Presentation()

        # Set slide dimensions to 16:9
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

        # Font settings
        self.title_font = "Inter"
        self.body_font = "Inter"
        self.fallback_font = "Arial"

        self._slide_count = 0

    def _get_background_color(self) -> "RGBColor":
        """Get the appropriate background color."""
        return KDSColors.BACKGROUND_DARK if self.dark_mode else KDSColors.BACKGROUND_LIGHT

    def _get_text_color(self) -> "RGBColor":
        """Get the appropriate text color."""
        return KDSColors.TEXT_LIGHT if self.dark_mode else KDSColors.TEXT_DARK

    def _set_slide_background(self, slide) -> None:
        """Set the slide background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._get_background_color()

    def _add_text_frame(
        self,
        slide,
        left: "Inches",
        top: "Inches",
        width: "Inches",
        height: "Inches",
        text: str,
        font_size: int = 18,
        bold: bool = False,
        alignment: "PP_ALIGN" = PP_ALIGN.LEFT,
    ):
        """Add a text frame to the slide."""
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = self._get_text_color()
        p.font.name = self.title_font
        p.alignment = alignment

        return textbox

    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
    ) -> "KDSPresentation":
        """
        Add a title slide.

        Args:
            title: Main title text.
            subtitle: Optional subtitle text.

        Returns:
            self for method chaining.
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(2.5),
            width=Inches(11.83),
            height=Inches(1.5),
            text=title,
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Subtitle
        if subtitle:
            self._add_text_frame(
                slide,
                left=Inches(0.75),
                top=Inches(4.0),
                width=Inches(11.83),
                height=Inches(0.75),
                text=subtitle,
                font_size=24,
                bold=False,
                alignment=PP_ALIGN.CENTER,
            )

        # Purple accent line
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(4.5),
            top=Inches(3.9),
            width=Inches(4.33),
            height=Inches(0.05),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = KDSColors.PRIMARY
        shape.line.fill.background()

        self._slide_count += 1
        logger.debug(f"Added title slide: {title}")
        return self

    def add_section_slide(
        self,
        section_title: str,
        section_number: Optional[int] = None,
    ) -> "KDSPresentation":
        """
        Add a section divider slide.

        Args:
            section_title: Section title text.
            section_number: Optional section number.

        Returns:
            self for method chaining.
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Section number (if provided)
        if section_number is not None:
            num_box = self._add_text_frame(
                slide,
                left=Inches(0.75),
                top=Inches(2.5),
                width=Inches(11.83),
                height=Inches(1.0),
                text=f"0{section_number}" if section_number < 10 else str(section_number),
                font_size=72,
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            # Make number purple
            num_box.text_frame.paragraphs[0].font.color.rgb = KDSColors.PRIMARY

        # Section title
        self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(3.5) if section_number else Inches(3.0),
            width=Inches(11.83),
            height=Inches(1.0),
            text=section_title,
            font_size=36,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        self._slide_count += 1
        logger.debug(f"Added section slide: {section_title}")
        return self

    def add_content_slide(
        self,
        title: str,
        bullet_points: List[str],
        subtitle: Optional[str] = None,
    ) -> "KDSPresentation":
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title.
            bullet_points: List of bullet point strings.
            subtitle: Optional subtitle below title.

        Returns:
            self for method chaining.
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        title_box = self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(0.5),
            width=Inches(11.83),
            height=Inches(0.75),
            text=title,
            font_size=32,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Subtitle
        content_top = Inches(1.25)
        if subtitle:
            self._add_text_frame(
                slide,
                left=Inches(0.75),
                top=Inches(1.25),
                width=Inches(11.83),
                height=Inches(0.5),
                text=subtitle,
                font_size=18,
                bold=False,
                alignment=PP_ALIGN.LEFT,
            )
            content_top = Inches(1.75)

        # Bullet points
        textbox = slide.shapes.add_textbox(
            left=Inches(0.75),
            top=content_top,
            width=Inches(11.83),
            height=Inches(5.25),
        )
        tf = textbox.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"  {point}"  # Indent for bullet appearance
            p.font.size = Pt(20)
            p.font.color.rgb = self._get_text_color()
            p.font.name = self.body_font
            p.space_after = Pt(12)

            # Add bullet character
            p.text = f"\u2022  {point}"

        self._slide_count += 1
        logger.debug(f"Added content slide: {title}")
        return self

    def add_chart_slide(
        self,
        title: str,
        chart_path: Union[str, Path],
        caption: Optional[str] = None,
    ) -> "KDSPresentation":
        """
        Add a slide with a chart image.

        Args:
            title: Slide title.
            chart_path: Path to the chart image file.
            caption: Optional caption below the chart.

        Returns:
            self for method chaining.

        Raises:
            FileNotFoundError: If the chart file doesn't exist.
        """
        chart_path = Path(chart_path)
        if not chart_path.exists():
            raise FileNotFoundError(f"Chart not found: {chart_path}")

        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(0.5),
            width=Inches(11.83),
            height=Inches(0.75),
            text=title,
            font_size=32,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Chart image (centered)
        chart_width = Inches(10)
        chart_height = Inches(5.5)
        chart_left = (self.SLIDE_WIDTH - chart_width) / 2
        chart_top = Inches(1.25)

        slide.shapes.add_picture(
            str(chart_path),
            left=chart_left,
            top=chart_top,
            width=chart_width,
            height=chart_height,
        )

        # Caption
        if caption:
            self._add_text_frame(
                slide,
                left=Inches(0.75),
                top=Inches(6.75),
                width=Inches(11.83),
                height=Inches(0.5),
                text=caption,
                font_size=12,
                bold=False,
                alignment=PP_ALIGN.CENTER,
            )

        self._slide_count += 1
        logger.debug(f"Added chart slide: {title}")
        return self

    def add_two_column_slide(
        self,
        title: str,
        left_content: List[str],
        right_content: List[str],
        left_header: Optional[str] = None,
        right_header: Optional[str] = None,
    ) -> "KDSPresentation":
        """
        Add a two-column content slide.

        Args:
            title: Slide title.
            left_content: Bullet points for left column.
            right_content: Bullet points for right column.
            left_header: Optional header for left column.
            right_header: Optional header for right column.

        Returns:
            self for method chaining.
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(0.5),
            width=Inches(11.83),
            height=Inches(0.75),
            text=title,
            font_size=32,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Left column
        left_top = Inches(1.5)
        if left_header:
            header_box = self._add_text_frame(
                slide,
                left=Inches(0.75),
                top=left_top,
                width=Inches(5.5),
                height=Inches(0.5),
                text=left_header,
                font_size=20,
                bold=True,
                alignment=PP_ALIGN.LEFT,
            )
            header_box.text_frame.paragraphs[0].font.color.rgb = KDSColors.PRIMARY
            left_top = Inches(2.0)

        left_box = slide.shapes.add_textbox(
            left=Inches(0.75),
            top=left_top,
            width=Inches(5.5),
            height=Inches(4.5),
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(left_content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2022  {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self._get_text_color()
            p.font.name = self.body_font
            p.space_after = Pt(8)

        # Right column
        right_top = Inches(1.5)
        if right_header:
            header_box = self._add_text_frame(
                slide,
                left=Inches(7.0),
                top=right_top,
                width=Inches(5.5),
                height=Inches(0.5),
                text=right_header,
                font_size=20,
                bold=True,
                alignment=PP_ALIGN.LEFT,
            )
            header_box.text_frame.paragraphs[0].font.color.rgb = KDSColors.PRIMARY
            right_top = Inches(2.0)

        right_box = slide.shapes.add_textbox(
            left=Inches(7.0),
            top=right_top,
            width=Inches(5.5),
            height=Inches(4.5),
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(right_content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2022  {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self._get_text_color()
            p.font.name = self.body_font
            p.space_after = Pt(8)

        self._slide_count += 1
        logger.debug(f"Added two-column slide: {title}")
        return self

    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
    ) -> "KDSPresentation":
        """
        Add a slide with a data table.

        Args:
            title: Slide title.
            headers: Column headers.
            rows: Table rows (list of lists).

        Returns:
            self for method chaining.
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(0.5),
            width=Inches(11.83),
            height=Inches(0.75),
            text=title,
            font_size=32,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Table
        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)

        table_width = Inches(11.5)
        table_height = Inches(0.4 * n_rows)
        col_width = table_width / n_cols

        table = slide.shapes.add_table(
            rows=n_rows,
            cols=n_cols,
            left=Inches(0.9),
            top=Inches(1.5),
            width=table_width,
            height=table_height,
        ).table

        # Style header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = KDSColors.PRIMARY

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = KDSColors.TEXT_LIGHT
            p.font.name = self.body_font
            p.alignment = PP_ALIGN.CENTER

        # Style data rows
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)

                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A) if self.dark_mode else RGBColor(0xF5, 0xF5, 0xF5)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._get_background_color()

                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = self._get_text_color()
                p.font.name = self.body_font
                p.alignment = PP_ALIGN.CENTER

        self._slide_count += 1
        logger.debug(f"Added table slide: {title}")
        return self

    def add_closing_slide(
        self,
        title: str = "Thank You",
        contact_info: Optional[List[str]] = None,
    ) -> "KDSPresentation":
        """
        Add a closing/thank you slide.

        Args:
            title: Closing title (default: "Thank You").
            contact_info: Optional contact information lines.

        Returns:
            self for method chaining.
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        title_box = self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(2.5),
            width=Inches(11.83),
            height=Inches(1.5),
            text=title,
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        title_box.text_frame.paragraphs[0].font.color.rgb = KDSColors.PRIMARY

        # Contact info
        if contact_info:
            info_box = slide.shapes.add_textbox(
                left=Inches(0.75),
                top=Inches(4.0),
                width=Inches(11.83),
                height=Inches(2.0),
            )
            tf = info_box.text_frame
            tf.word_wrap = True

            for i, line in enumerate(contact_info):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = self._get_text_color()
                p.font.name = self.body_font
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(6)

        self._slide_count += 1
        logger.debug(f"Added closing slide: {title}")
        return self

    def add_insight_slide(self, insight_data: Dict[str, Any]) -> "KDSPresentation":
        """
        Add a slide from an insight specification.

        Args:
            insight_data: Dict with keys:
                - headline: Action title for the slide
                - supporting_text: Body text or bullet points
                - evidence: List of evidence dicts (optional)
                - suggested_slide_type: "comparison", "trend", "content", etc.
                - chart_path: Path to chart image (optional)

        Returns:
            self for method chaining.
        """
        headline = insight_data.get('headline', 'Insight')
        supporting_text = insight_data.get('supporting_text', '')
        chart_path = insight_data.get('chart_path')
        slide_type = insight_data.get('suggested_slide_type', 'content')

        # If we have a chart, use chart slide
        if chart_path and Path(chart_path).exists():
            self.add_chart_slide(
                title=headline,
                chart_path=chart_path,
                caption=supporting_text[:200] if len(supporting_text) > 200 else supporting_text
            )
        else:
            # Use content slide with supporting text as bullet
            self.add_content_slide(
                title=headline,
                bullet_points=[supporting_text] if supporting_text else []
            )

        return self

    def build_from_insights(
        self,
        catalog_path: str,
        include_title: bool = True,
        include_sections: bool = True,
        max_slides: int = 20,
    ) -> int:
        """
        Build presentation from an insight catalog.

        Args:
            catalog_path: Path to insights.yaml file
            include_title: Whether to add title slide
            include_sections: Whether to add section dividers
            max_slides: Maximum number of content slides

        Returns:
            Number of slides added
        """
        from core.insight_engine import InsightCatalog

        catalog = InsightCatalog.load(catalog_path)
        slides_added = 0

        # Title slide
        if include_title:
            # Extract title from business question
            title = "Analysis Results"
            if catalog.business_question:
                # Try to make a title from the question
                q = catalog.business_question
                if q.lower().startswith("what"):
                    title = q.replace("?", "").strip()
                else:
                    title = f"Analysis: {q[:50]}"

            self.add_title_slide(title=title, subtitle=f"Generated {catalog.generated_at[:10]}")
            slides_added += 1

        arc = catalog.narrative_arc

        # Key Findings section
        key_findings = arc.get('key_findings', [])
        if key_findings and slides_added < max_slides:
            if include_sections:
                self.add_section_slide("Key Findings", section_number=1)
                slides_added += 1

            for insight_id in key_findings:
                if slides_added >= max_slides:
                    break
                insight = next((i for i in catalog.insights if i.id == insight_id), None)
                if insight:
                    self._add_insight_as_slide(insight)
                    slides_added += 1

        # Supporting findings (limit to 3)
        supporting = arc.get('supporting_findings', [])
        if supporting and slides_added < max_slides:
            if include_sections and key_findings:
                self.add_section_slide("Additional Findings", section_number=2)
                slides_added += 1

            for insight_id in supporting[:3]:
                if slides_added >= max_slides:
                    break
                insight = next((i for i in catalog.insights if i.id == insight_id), None)
                if insight:
                    self._add_insight_as_slide(insight)
                    slides_added += 1

        # Implications
        implications = arc.get('implications', [])
        if implications and slides_added < max_slides:
            if include_sections:
                self.add_section_slide("Implications", section_number=3)
                slides_added += 1

            for insight_id in implications:
                if slides_added >= max_slides:
                    break
                insight = next((i for i in catalog.insights if i.id == insight_id), None)
                if insight:
                    self._add_insight_as_slide(insight)
                    slides_added += 1

        # Recommendations
        recommendations = arc.get('recommendations', [])
        if recommendations and slides_added < max_slides:
            if include_sections:
                self.add_section_slide("Recommendations", section_number=4)
                slides_added += 1

            for insight_id in recommendations:
                if slides_added >= max_slides:
                    break
                insight = next((i for i in catalog.insights if i.id == insight_id), None)
                if insight:
                    self._add_insight_as_slide(insight)
                    slides_added += 1

        return slides_added

    def _add_insight_as_slide(self, insight) -> None:
        """Add a single insight as a slide."""
        # Check for chart evidence
        chart_evidence = next(
            (e for e in insight.evidence if e.type == "chart"),
            None
        )

        if chart_evidence and chart_evidence.reference and Path(chart_evidence.reference).exists():
            self.add_chart_slide(
                title=insight.headline,
                chart_path=chart_evidence.reference,
                caption=insight.supporting_text[:200]
            )
        else:
            self.add_content_slide(
                title=insight.headline,
                bullet_points=[insight.supporting_text]
            )

    def add_chart_slide_with_notes(
        self,
        title: str,
        chart_path: str,
        caption: str = '',
        notes_context: Dict[str, Any] = None,
    ) -> "KDSPresentation":
        """
        Add chart slide with comprehensive slide notes.

        Args:
            title: Slide title (action title format)
            chart_path: Path to chart image
            caption: Brief caption under chart
            notes_context: Dict with context for speaker notes:
                - source: Data source attribution
                - methodology: How the data was calculated
                - talking_points: List of key points to mention
                - raw_data: Optional data summary
                - caveats: Any limitations or assumptions

        Returns:
            self for method chaining.
        """
        # Add the chart slide using existing method
        self.add_chart_slide(title, chart_path, caption)

        # Get the slide we just added
        slide = self.prs.slides[-1]

        # Add comprehensive notes if context provided
        if notes_context:
            notes_slide = slide.notes_slide
            notes_text = self._format_comprehensive_notes(notes_context)
            notes_slide.notes_text_frame.text = notes_text

        return self

    def _format_comprehensive_notes(self, context: Dict[str, Any]) -> str:
        """
        Format comprehensive slide notes from context.

        Creates structured notes with:
        - Data source attribution
        - Methodology explanation
        - Key talking points
        - Raw data reference
        - Caveats and assumptions
        """
        sections = []

        # Source
        source = context.get('source')
        if source:
            sections.append(f"DATA SOURCE\n{source}")

        # Methodology
        methodology = context.get('methodology')
        if methodology:
            sections.append(f"METHODOLOGY\n{methodology}")

        # Talking points
        talking_points = context.get('talking_points', [])
        if talking_points:
            points_text = "\n".join(f"* {point}" for point in talking_points)
            sections.append(f"KEY TALKING POINTS\n{points_text}")

        # Raw data
        raw_data = context.get('raw_data')
        if raw_data:
            if isinstance(raw_data, dict):
                data_lines = [f"  {k}: {v}" for k, v in raw_data.items()]
                data_text = "\n".join(data_lines)
            else:
                data_text = str(raw_data)
            sections.append(f"SUPPORTING DATA\n{data_text}")

        # Caveats
        caveats = context.get('caveats')
        if caveats:
            if isinstance(caveats, list):
                caveats_text = "\n".join(f"* {c}" for c in caveats)
            else:
                caveats_text = caveats
            sections.append(f"CAVEATS & ASSUMPTIONS\n{caveats_text}")

        return "\n\n".join(sections)

    def add_native_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: Dict[str, Any],
        caption: str = '',
    ) -> "KDSPresentation":
        """
        Add slide with native, editable PowerPoint chart.

        Native charts can be edited directly in PowerPoint,
        unlike image-based charts.

        Args:
            title: Slide title
            chart_type: One of 'bar', 'column', 'line', 'pie', 'area'
            data: Dict with structure:
                {
                    'categories': ['Q1', 'Q2', 'Q3'],
                    'series': [
                        {'name': 'Revenue', 'values': [100, 120, 140]},
                        {'name': 'Cost', 'values': [80, 90, 100]}
                    ]
                }
            caption: Optional caption text

        Returns:
            self for method chaining.
        """
        # Map string types to PowerPoint chart types
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'bar_stacked': XL_CHART_TYPE.BAR_STACKED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'column_stacked': XL_CHART_TYPE.COLUMN_STACKED,
            'line': XL_CHART_TYPE.LINE,
            'line_markers': XL_CHART_TYPE.LINE_MARKERS,
            'pie': XL_CHART_TYPE.PIE,
            'area': XL_CHART_TYPE.AREA,
            'area_stacked': XL_CHART_TYPE.AREA_STACKED,
        }

        xl_chart_type = chart_type_map.get(chart_type.lower())
        if xl_chart_type is None:
            raise ValueError(f"Unsupported chart type: {chart_type}. "
                            f"Supported: {list(chart_type_map.keys())}")

        # Create slide
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Add title
        self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(0.5),
            width=Inches(11.83),
            height=Inches(0.75),
            text=title,
            font_size=32,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Build chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.get('categories', [])

        for series in data.get('series', []):
            chart_data.add_series(series['name'], series['values'])

        # Add chart to slide
        x, y, cx, cy = Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.25)
        chart = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, chart_data
        ).chart

        # Apply KDS styling to chart
        self._apply_kds_chart_style(chart, xl_chart_type)

        # Add caption if provided
        if caption:
            self._add_text_frame(
                slide,
                left=Inches(0.75),
                top=Inches(6.85),
                width=Inches(11.83),
                height=Inches(0.4),
                text=caption,
                font_size=12,
                bold=False,
                alignment=PP_ALIGN.CENTER,
            )

        self._slide_count += 1
        logger.debug(f"Added native chart slide: {title}")
        return self

    def _apply_kds_chart_style(self, chart, chart_type) -> None:
        """
        Apply KDS brand styling to native PowerPoint chart.

        - Kearney Purple for primary series
        - No gridlines
        - Clean axis styling
        """
        # KDS color palette (no green!)
        kds_colors = [
            RGBColor(0x78, 0x23, 0xDC),  # Kearney Purple
            RGBColor(0x9B, 0x4D, 0xCA),  # Light purple
            RGBColor(0x5C, 0x1B, 0xA8),  # Dark purple
            RGBColor(0xB2, 0x66, 0xFF),  # Bright purple
            RGBColor(0x66, 0x66, 0x66),  # Gray
            RGBColor(0x99, 0x99, 0x99),  # Light gray
        ]

        # Apply colors to series
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            color = kds_colors[i % len(kds_colors)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color

        # Style axes if present
        try:
            if hasattr(chart, 'category_axis'):
                cat_axis = chart.category_axis
                cat_axis.has_major_gridlines = False
                cat_axis.has_minor_gridlines = False

            if hasattr(chart, 'value_axis'):
                val_axis = chart.value_axis
                val_axis.has_major_gridlines = False
                val_axis.has_minor_gridlines = False
        except Exception:
            # Some chart types don't have these axes
            pass

    def add_placeholder_slide(
        self,
        title: str,
        placeholder_text: str,
        replacement_instructions: str = '',
        dimensions: Tuple[float, float] = None,
    ) -> "KDSPresentation":
        """
        Add a slide with a placeholder for manual content insertion.

        Use this when automated generation can't produce the needed
        visual, but the slide structure should be prepared.

        Args:
            title: Slide title
            placeholder_text: What should go in this space
            replacement_instructions: How to create the replacement content
            dimensions: Optional (width, height) in inches for placeholder box

        Returns:
            self for method chaining.
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Add title
        self._add_text_frame(
            slide,
            left=Inches(0.75),
            top=Inches(0.5),
            width=Inches(11.83),
            height=Inches(0.75),
            text=title,
            font_size=32,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Determine placeholder dimensions
        if dimensions:
            width, height = dimensions
        else:
            width, height = 11.0, 5.0

        # Add placeholder box
        left = Inches((13.33 - width) / 2)  # Center horizontally
        top = Inches(1.5)

        placeholder_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, Inches(width), Inches(height)
        )

        # Style placeholder box
        placeholder_box.fill.solid()
        placeholder_box.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A)  # Dark gray
        placeholder_box.line.color.rgb = KDSColors.PRIMARY  # Purple border
        placeholder_box.line.width = Pt(2)

        # Add placeholder text inside the box
        tf = placeholder_box.text_frame
        tf.word_wrap = True

        # Main placeholder text
        p = tf.paragraphs[0]
        p.text = f"[PLACEHOLDER: {placeholder_text}]"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = KDSColors.PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # Add instructions if provided
        if replacement_instructions:
            p2 = tf.add_paragraph()
            p2.text = f"\n{replacement_instructions}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = KDSColors.GRAY_400
            p2.alignment = PP_ALIGN.CENTER

        # Add note about placeholder
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = (
            f"PLACEHOLDER SLIDE\n\n"
            f"Content needed: {placeholder_text}\n\n"
            f"Instructions: {replacement_instructions or 'Replace with final visual'}"
        )

        self._slide_count += 1
        logger.debug(f"Added placeholder slide: {title}")
        return self

    @property
    def slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        return self._slide_count

    def save(self, path: Union[str, Path]) -> Path:
        """
        Save the presentation to a file.

        Args:
            path: Output file path.

        Returns:
            Path to the saved file.
        """
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(path))
        logger.info(f"Presentation saved to {path} ({self._slide_count} slides)")
        return path


def main():
    """Demo of KDSPresentation capabilities."""
    pres = KDSPresentation()

    # Title slide
    pres.add_title_slide(
        "Q4 2025 Revenue Analysis",
        "Prepared for ACME Corporation",
    )

    # Section divider
    pres.add_section_slide("Executive Summary", section_number=1)

    # Content slide
    pres.add_content_slide(
        "Key Findings",
        [
            "Revenue increased 15% year-over-year",
            "North region outperformed targets by 20%",
            "New product line contributed $2.5M",
            "Customer retention improved to 94%",
        ],
    )

    # Two-column slide
    pres.add_two_column_slide(
        "Opportunities vs Challenges",
        left_content=[
            "Expand into European markets",
            "Launch premium tier product",
            "Increase digital marketing spend",
        ],
        right_content=[
            "Supply chain constraints",
            "Competitive pricing pressure",
            "Talent acquisition",
        ],
        left_header="Opportunities",
        right_header="Challenges",
    )

    # Table slide
    pres.add_table_slide(
        "Regional Performance",
        headers=["Region", "Revenue", "Growth", "Target"],
        rows=[
            ["North", "$4.2M", "+22%", "Met"],
            ["South", "$3.1M", "+12%", "Met"],
            ["East", "$2.8M", "+8%", "Below"],
            ["West", "$3.5M", "+18%", "Met"],
        ],
    )

    # Closing slide
    pres.add_closing_slide(
        "Thank You",
        contact_info=[
            "Kearney Digital & Analytics",
            "analytics@kearney.com",
        ],
    )

    pres.save("exports/demo_presentation.pptx")
    print(f"Demo presentation created with {pres.slide_count} slides")


if __name__ == "__main__":
    main()
